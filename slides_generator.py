"""
Professional PPTX slayd generator v2
- AI mavzuga qarab rang palitrasini tanlaydi
- Har bir slayd orqa foni AI yaratgan RASM bilan
- Har bir varaq HAR XIL dizaynda (6 xil layout aylanadi)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image, ImageDraw, ImageFilter, ImageEnhance
import io
import re
import requests
import urllib.parse
import random
import logging
from concurrent.futures import ThreadPoolExecutor

log = logging.getLogger("slidebot")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

DEFAULT_THEME = {
    "primary": RGBColor(0x0D, 0x47, 0xA1),
    "accent": RGBColor(0x00, 0xC2, 0xE0),
    "dark": RGBColor(0x0A, 0x18, 0x2E),
    "light": RGBColor(0xEA, 0xF3, 0xFF),
}

# ===== 10 TA FAN DIZAYNI (foydalanuvchi tanlaydi) =====
SUBJECT_THEMES = {
    "matematika": {"primary": "#123C69", "accent": "#F2A104", "dark": "#0A1F33", "light": "#EAF1F8"},
    "fizika":     {"primary": "#1A237E", "accent": "#00E5FF", "dark": "#0D1240", "light": "#E8EAF6"},
    "kimyo":      {"primary": "#00695C", "accent": "#AEEA00", "dark": "#00332C", "light": "#E0F2F1"},
    "biologiya":  {"primary": "#2E7D32", "accent": "#FFC107", "dark": "#14380F", "light": "#E8F5E9"},
    "tarix":      {"primary": "#6D4C41", "accent": "#D4AF37", "dark": "#2E1A12", "light": "#EFEBE9"},
    "adabiyot":   {"primary": "#880E4F", "accent": "#FF8A65", "dark": "#3D0625", "light": "#FCE4EC"},
    "geografiya": {"primary": "#0277BD", "accent": "#66BB6A", "dark": "#012A42", "light": "#E1F5FE"},
    "informatika":{"primary": "#37474F", "accent": "#00E676", "dark": "#10171B", "light": "#ECEFF1"},
    "iqtisodiyot":{"primary": "#1B2A4A", "accent": "#2ECC71", "dark": "#0E1730", "light": "#EDF2F9"},
    "ingliz":     {"primary": "#B71C1C", "accent": "#1565C0", "dark": "#3C0A0A", "light": "#FFEBEE"},
}

# Har fan dizayni uchun 4 ta stok rasm — Mini App preview'dagi bilan BIR XIL!
_PX = "https://images.pexels.com/photos/{0}/pexels-photo-{0}.{1}?auto=compress&cs=tinysrgb&fit=crop&h=720&w=1280"
SUBJECT_IMAGES = {
    "matematika": [_PX.format(6238050,'jpeg'), _PX.format(6256066,'jpeg'), _PX.format(5412336,'jpeg'), _PX.format(5412287,'jpeg')],
    "fizika":     [_PX.format(13014236,'jpeg'), _PX.format(27088485,'jpeg'), _PX.format(7723393,'jpeg'), _PX.format(37269536,'jpeg')],
    "kimyo":      [_PX.format(8533087,'jpeg'), _PX.format(6208943,'jpeg'), _PX.format(5427863,'jpeg'), _PX.format(8533140,'jpeg')],
    "biologiya":  [_PX.format(18069423,'png'), _PX.format(1302883,'jpeg'), _PX.format(10448365,'jpeg'), _PX.format(12642803,'jpeg')],
    "tarix":      [_PX.format(35284819,'jpeg'), _PX.format(19227934,'jpeg'), _PX.format(6493636,'jpeg'), _PX.format(16386337,'jpeg')],
    "adabiyot":   [_PX.format(13923801,'jpeg'), _PX.format(34640246,'jpeg'), _PX.format(13921129,'jpeg'), _PX.format(34412232,'jpeg')],
    "geografiya": [_PX.format(1098518,'jpeg'), _PX.format(23848552,'jpeg'), _PX.format(7412075,'jpeg'), _PX.format(18383537,'jpeg')],
    "informatika":[_PX.format(6424583,'jpeg'), _PX.format(5480781,'jpeg'), _PX.format(7325498,'jpeg'), _PX.format(17489157,'jpeg')],
    "iqtisodiyot":[_PX.format(186461,'jpeg'), _PX.format(6801639,'jpeg'), _PX.format(7792841,'jpeg'), _PX.format(5398879,'jpeg')],
    "ingliz":     [_PX.format(13887712,'jpeg'), _PX.format(7450476,'jpeg'), _PX.format(17379169,'jpeg'), _PX.format(4440715,'jpeg')],
}


def _fetch_url_image(url: str) -> bytes | None:
    """URL'dan rasm yuklab olish"""
    try:
        resp = requests.get(url, timeout=20)
        if resp.status_code == 200 and len(resp.content) > 5000:
            return resp.content
    except Exception:
        pass
    return None


def hex_to_rgb(hex_str: str) -> RGBColor:
    """#RRGGBB -> RGBColor"""
    try:
        h = hex_str.lstrip('#')
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        return DEFAULT_THEME["primary"]


def build_theme(theme_colors: dict | None) -> dict:
    """AI tanlagan ranglardan tema yasash"""
    if not theme_colors:
        return dict(DEFAULT_THEME)
    theme = dict(DEFAULT_THEME)
    if theme_colors.get("primary"):
        theme["primary"] = hex_to_rgb(theme_colors["primary"])
    if theme_colors.get("accent"):
        theme["accent"] = hex_to_rgb(theme_colors["accent"])
    if theme_colors.get("dark"):
        theme["dark"] = hex_to_rgb(theme_colors["dark"])
    if theme_colors.get("light"):
        theme["light"] = hex_to_rgb(theme_colors["light"])
    return theme


# ==================== RASM YUKLASH ====================

def _is_image(resp) -> bool:
    """Javob haqiqiy rasmmi?"""
    ct = resp.headers.get('content-type', '')
    return resp.status_code == 200 and len(resp.content) > 5000 and 'image' in ct


def _fetch_stock_photo(keywords: str, width: int = 1024, height: int = 576, seed: int = 0) -> bytes | None:
    """LoremFlickr — HAQIQIY stok fotolar (juda ishonchli va tez zaxira)"""
    # kalit so'zlarni tozalash: faqat oddiy inglizcha so'zlar
    words = re.findall(r'[a-zA-Z]+', keywords)[:3]
    if not words:
        words = ["business"]
    tags = ",".join(words)
    try:
        url = f"https://loremflickr.com/{width}/{height}/{urllib.parse.quote(tags)}?lock={seed % 100}"
        resp = requests.get(url, timeout=15)
        if _is_image(resp):
            return resp.content
    except Exception:
        pass
    return None


def _fetch_ai_image(prompt: str, width: int = 1024, height: int = 576, seed: int = 0,
                    fast: bool = False) -> bytes | None:
    """Rasm olish: Pollinations AI (flux) → turbo → LoremFlickr (real foto).
    Har doim rasm qaytarishga harakat qiladi! fast=True — flux o'tkazib yuboriladi"""
    # 1-urinish: flux (eng sifatli AI rasm) — fast rejimda o'tkazib yuboriladi
    if not fast:
        try:
            enc = urllib.parse.quote(prompt[:200])
            url = (f"https://image.pollinations.ai/prompt/{enc}"
                   f"?width={width}&height={height}&nologo=true&seed={seed}&model=flux")
            resp = requests.get(url, timeout=40)
            if _is_image(resp):
                return resp.content
        except Exception:
            pass
    # 2-urinish: turbo (tez AI rasm)
    try:
        enc = urllib.parse.quote(prompt[:150])
        url = (f"https://image.pollinations.ai/prompt/{enc}"
               f"?width={width}&height={height}&nologo=true&seed={seed}&model=turbo")
        resp = requests.get(url, timeout=25)
        if _is_image(resp):
            return resp.content
    except Exception:
        pass
    # Stok foto ishlatmaymiz — mavzuga aloqasiz rasm chiqib qoladi
    return None


def _fallback_gradient(width: int, height: int, theme: dict) -> bytes:
    """Rasm yuklanmasa — chiroyli gradient fon"""
    img = Image.new('RGB', (width, height))
    draw = ImageDraw.Draw(img)
    c1 = (theme["dark"][0], theme["dark"][1], theme["dark"][2])
    c2 = (theme["primary"][0], theme["primary"][1], theme["primary"][2])
    for y in range(height):
        r = y / height
        draw.line([(0, y), (width, y)], fill=(
            int(c1[0] * (1 - r) + c2[0] * r),
            int(c1[1] * (1 - r) + c2[1] * r),
            int(c1[2] * (1 - r) + c2[2] * r),
        ))
    # dekorativ doiralar
    rnd = random.Random(42)
    overlay = Image.new('RGBA', (width, height), (0, 0, 0, 0))
    od = ImageDraw.Draw(overlay)
    ac = (theme["accent"][0], theme["accent"][1], theme["accent"][2])
    for _ in range(6):
        x, y = rnd.randint(0, width), rnd.randint(0, height)
        rr = rnd.randint(60, 220)
        od.ellipse([x - rr, y - rr, x + rr, y + rr], fill=(ac[0], ac[1], ac[2], 25))
    img = Image.alpha_composite(img.convert('RGBA'), overlay).convert('RGB')
    buf = io.BytesIO()
    img.save(buf, format='JPEG', quality=88)
    return buf.getvalue()


def _prepare_image(raw: bytes, width: int, height: int,
                   darken: float = 0.0, blur: float = 0.0) -> bytes:
    """Rasmni kerakli o'lchamga moslash (crop), qoraytirish va blur"""
    img = Image.open(io.BytesIO(raw)).convert('RGB')

    # Crop to aspect
    target_ratio = width / height
    w, h = img.size
    ratio = w / h
    if ratio > target_ratio:
        new_w = int(h * target_ratio)
        left = (w - new_w) // 2
        img = img.crop((left, 0, left + new_w, h))
    else:
        new_h = int(w / target_ratio)
        top = (h - new_h) // 2
        img = img.crop((0, top, w, top + new_h))
    img = img.resize((width, height), Image.LANCZOS)

    if blur > 0:
        img = img.filter(ImageFilter.GaussianBlur(radius=blur))
    if darken > 0:
        img = ImageEnhance.Brightness(img).enhance(1 - darken)

    buf = io.BytesIO()
    img.save(buf, format='JPEG', quality=88)
    return buf.getvalue()


def _rounded_image(raw: bytes, width: int, height: int, radius: int = 40) -> bytes:
    """Burchagi yumaloq rasm (PNG, shaffof burchak)"""
    prepared = _prepare_image(raw, width, height)
    img = Image.open(io.BytesIO(prepared)).convert('RGBA')
    mask = Image.new('L', (width, height), 0)
    md = ImageDraw.Draw(mask)
    md.rounded_rectangle([0, 0, width, height], radius=radius, fill=255)
    img.putalpha(mask)
    buf = io.BytesIO()
    img.save(buf, format='PNG')
    return buf.getvalue()


# ==================== YORDAMCHI ====================

def _set_font(run, size=None, bold=None, color=None, font_name="Arial"):
    if size:
        run.font.size = size
    if bold is not None:
        run.font.bold = bold
    if color:
        run.font.color.rgb = color
    run.font.name = font_name
    rPr = run._r.get_or_add_rPr()
    rPr.set(qn('a:latin'), font_name)


def _text(slide, left, top, width, height, text, size=Pt(16), color=None,
          bold=False, align=PP_ALIGN.LEFT, anchor='t', line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    _set_font(run, size=size, bold=bold, color=color)
    tf._txBody.bodyPr.set('anchor', anchor)
    return box


def _rect(slide, left, top, width, height, color, transparency: int = 0):
    """Rangli to'rtburchak. transparency: 0-100 (%)"""
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    if transparency > 0:
        # XML orqali shaffoflik
        sPr = sh.fill._xPr.find(qn('a:solidFill'))
        srgb = sPr.find(qn('a:srgbClr'))
        alpha = srgb.makeelement(qn('a:alpha'), {'val': str(int((100 - transparency) * 1000))})
        srgb.append(alpha)
    return sh


def _full_bg_image(slide, img_bytes: bytes):
    """Slaydga to'liq fon rasmi"""
    slide.shapes.add_picture(io.BytesIO(img_bytes), 0, 0, SLIDE_W, SLIDE_H)


WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ==================== LAYOUTLAR ====================

def _layout_title(prs, slide_item, theme, img_raw, pres_title):
    """1. SARLAVHA: to'liq rasm fon + qoraytirilgan + katta sarlavha"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = _prepare_image(img_raw, 1280, 720, darken=0.55, blur=1.5)
    _full_bg_image(slide, bg)

    # aksent chiziq
    _rect(slide, Inches(1), Inches(2.1), Inches(1.6), Pt(6), theme["accent"])

    _text(slide, Inches(1), Inches(2.4), Inches(11), Inches(2.2),
          slide_item.get('title', pres_title), size=Pt(48), bold=True, color=WHITE)

    sub = slide_item.get('subtitle') or slide_item.get('content', '')[:180]
    if sub:
        _text(slide, Inches(1), Inches(4.7), Inches(10.5), Inches(1.5),
              sub, size=Pt(20), color=RGBColor(0xE8, 0xE8, 0xE8))

    # Muallif ismi (pastda)
    author = slide_item.get('author', '')
    if author:
        _rect(slide, Inches(1), Inches(6.55), Inches(0.08), Inches(0.5), theme["accent"])
        _text(slide, Inches(1.25), Inches(6.55), Inches(9), Inches(0.5),
              f"Tayyorladi: {author}", size=Pt(16), bold=True, color=WHITE)


def _layout_image_right(prs, slide_item, theme, img_raw, num, img2_raw=None):
    """2. MATN CHAP + RASM O'NG (yumaloq burchakli), pastda 2-rasm"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    # chap tomonda yupqa rangli panel
    _rect(slide, 0, 0, Inches(0.25), SLIDE_H, theme["primary"])

    _text(slide, Inches(0.8), Inches(0.5), Inches(6.6), Inches(1.2),
          slide_item.get('title', ''), size=Pt(30), bold=True, color=theme["dark"])
    _rect(slide, Inches(0.8), Inches(1.55), Inches(2.2), Pt(4), theme["accent"])

    if img2_raw:
        # Matn qisqaroq joy oladi, pastda 2-rasm
        _text(slide, Inches(0.8), Inches(1.95), Inches(6.4), Inches(3.1),
              slide_item.get('content', ''), size=Pt(14), color=RGBColor(0x33, 0x33, 0x3D))
        img2 = _rounded_image(img2_raw, 640, 360, radius=30)
        slide.shapes.add_picture(io.BytesIO(img2), Inches(0.8), Inches(5.0), Inches(4.2), Inches(2.1))
    else:
        _text(slide, Inches(0.8), Inches(1.95), Inches(6.4), Inches(5.1),
              slide_item.get('content', ''), size=Pt(15), color=RGBColor(0x33, 0x33, 0x3D))

    img = _rounded_image(img_raw, 720, 800, radius=40)
    slide.shapes.add_picture(io.BytesIO(img), Inches(7.7), Inches(0.9), Inches(5.0), Inches(5.6))

    _text(slide, Inches(12.5), Inches(7.05), Inches(0.6), Inches(0.35),
          str(num), size=Pt(12), color=theme["primary"], align=PP_ALIGN.RIGHT)


def _layout_image_left(prs, slide_item, theme, img_raw, num, img2_raw=None):
    """3. RASM CHAP (to'liq balandlik) + MATN O'NG rangli fonda, 2-rasm pastda"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme["dark"])

    img = _prepare_image(img_raw, 640, 800)
    slide.shapes.add_picture(io.BytesIO(img), 0, 0, Inches(5.6), SLIDE_H)

    # rasm ustidan o'ngga o'tuvchi soyali panel
    _rect(slide, Inches(5.6), 0, Inches(0.12), SLIDE_H, theme["accent"])

    _text(slide, Inches(6.2), Inches(0.7), Inches(6.5), Inches(1.3),
          slide_item.get('title', ''), size=Pt(28), bold=True, color=WHITE)
    _rect(slide, Inches(6.2), Inches(1.85), Inches(1.8), Pt(4), theme["accent"])

    if img2_raw:
        _text(slide, Inches(6.2), Inches(2.25), Inches(6.5), Inches(2.6),
              slide_item.get('content', ''), size=Pt(14),
              color=RGBColor(0xE2, 0xE6, 0xEE))
        img2 = _rounded_image(img2_raw, 640, 340, radius=30)
        slide.shapes.add_picture(io.BytesIO(img2), Inches(6.2), Inches(5.0), Inches(4.4), Inches(2.2))
    else:
        _text(slide, Inches(6.2), Inches(2.25), Inches(6.5), Inches(4.6),
              slide_item.get('content', ''), size=Pt(15),
              color=RGBColor(0xE2, 0xE6, 0xEE))

    _text(slide, Inches(12.5), Inches(7.05), Inches(0.6), Inches(0.35),
          str(num), size=Pt(12), color=WHITE, align=PP_ALIGN.RIGHT)


def _layout_full_image(prs, slide_item, theme, img_raw, num):
    """4. TO'LIQ RASM FON + pastda matn paneli"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = _prepare_image(img_raw, 1280, 720, darken=0.25)
    _full_bg_image(slide, bg)

    # pastki qorong'i panel
    _rect(slide, 0, Inches(4.2), SLIDE_W, Inches(3.3), theme["dark"], transparency=18)

    _text(slide, Inches(0.9), Inches(4.5), Inches(11.5), Inches(0.9),
          slide_item.get('title', ''), size=Pt(28), bold=True, color=WHITE)
    _rect(slide, Inches(0.9), Inches(5.35), Inches(1.8), Pt(4), theme["accent"])
    _text(slide, Inches(0.9), Inches(5.6), Inches(11.5), Inches(1.7),
          slide_item.get('content', ''), size=Pt(14), color=RGBColor(0xE8, 0xE8, 0xEE))

    _text(slide, Inches(12.5), Inches(0.3), Inches(0.6), Inches(0.35),
          str(num), size=Pt(12), color=WHITE, align=PP_ALIGN.RIGHT)


def _layout_bullets(prs, slide_item, theme, img_raw, num):
    """5. NUQTALI KARTALAR: yuqorida rasm banner + pastda rangli kartalar"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0xF5, 0xF7, 0xFA))

    # yuqori banner rasm
    banner = _prepare_image(img_raw, 1280, 256, darken=0.45)
    slide.shapes.add_picture(io.BytesIO(banner), 0, 0, SLIDE_W, Inches(1.9))

    _text(slide, Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0),
          slide_item.get('title', ''), size=Pt(30), bold=True, color=WHITE)

    items = slide_item.get('bullet_points') or []
    if not items:
        items = [s.strip() for s in slide_item.get('content', '').split('.') if s.strip()][:6]
    items = [it.strip().lstrip('-•* ') for it in items if it.strip()]

    img2_raw = slide_item.get('_img2')
    if img2_raw:
        # O'ngda 2-rasm, kartalar chapda bitta ustun bo'lib
        items = items[:3]
        img2 = _rounded_image(img2_raw, 560, 600, radius=30)
        slide.shapes.add_picture(io.BytesIO(img2), Inches(8.6), Inches(2.25), Inches(4.2), Inches(4.5))
        card_w_override = Inches(7.6)
    else:
        items = items[:6]
        card_w_override = None

    # kartalar: 2 ustunli (yoki rasm bo'lsa 1 ustunli keng)
    top0 = Inches(2.25)
    card_w = card_w_override or Inches(6.0)
    card_h = Inches(1.5)
    gap = Inches(0.25)
    single_col = card_w_override is not None
    for i, item in enumerate(items):
        col = 0 if single_col else i % 2
        row = i if single_col else i // 2
        left = Inches(0.5) + col * (card_w + gap)
        top = top0 + row * (card_h + gap)
        card = _rect(slide, left, top, card_w, card_h, WHITE)
        card.line.color.rgb = RGBColor(0xE0, 0xE4, 0xEA)
        card.line.width = Pt(1)
        # raqamli aksent doira
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.18),
                                      top + Inches(0.45), Inches(0.55), Inches(0.55))
        circ.fill.solid()
        circ.fill.fore_color.rgb = theme["accent"] if i % 2 == 0 else theme["primary"]
        circ.line.fill.background()
        tfc = circ.text_frame
        pc = tfc.paragraphs[0]
        pc.alignment = PP_ALIGN.CENTER
        rc = pc.add_run()
        rc.text = str(i + 1)
        _set_font(rc, size=Pt(16), bold=True, color=WHITE)

        _text(slide, left + Inches(0.95), top + Inches(0.12), card_w - Inches(1.1),
              card_h - Inches(0.24), item, size=Pt(12.5),
              color=RGBColor(0x2A, 0x2A, 0x35), anchor='ctr')

    _text(slide, Inches(12.5), Inches(7.05), Inches(0.6), Inches(0.35),
          str(num), size=Pt(12), color=theme["primary"], align=PP_ALIGN.RIGHT)


def _layout_split_diagonal(prs, slide_item, theme, img_raw, num, img2_raw=None):
    """6. YUQORI RASM + PASTKI MATN + o'ngda 2-rasm"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)

    # yuqori yarmida keng rasm
    img = _prepare_image(img_raw, 1280, 448, darken=0.15)
    slide.shapes.add_picture(io.BytesIO(img), 0, 0, SLIDE_W, Inches(3.3))

    # rasm ustiga chapdan rangli sarlavha paneli
    _rect(slide, Inches(0.7), Inches(2.55), Inches(7.6), Inches(1.35), theme["primary"])
    _text(slide, Inches(1.0), Inches(2.7), Inches(7.0), Inches(1.05),
          slide_item.get('title', ''), size=Pt(24), bold=True, color=WHITE, anchor='ctr')

    if img2_raw:
        # Matn chapda, 2-rasm o'ngda
        _text(slide, Inches(0.8), Inches(4.25), Inches(7.6), Inches(2.9),
              slide_item.get('content', ''), size=Pt(14), color=RGBColor(0x33, 0x33, 0x3D))
        img2 = _rounded_image(img2_raw, 560, 420, radius=30)
        slide.shapes.add_picture(io.BytesIO(img2), Inches(8.8), Inches(4.2), Inches(3.9), Inches(2.9))
    else:
        _text(slide, Inches(0.8), Inches(4.25), Inches(11.7), Inches(2.9),
              slide_item.get('content', ''), size=Pt(15), color=RGBColor(0x33, 0x33, 0x3D))

    sub = slide_item.get('subtitle', '')
    if sub:
        _rect(slide, Inches(0.8), Inches(6.75), Inches(11.7), Pt(2), theme["accent"])
        _text(slide, Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.5),
              "💡 " + sub, size=Pt(12), color=theme["primary"])

    _text(slide, Inches(12.5), Inches(7.05), Inches(0.6), Inches(0.35),
          str(num), size=Pt(12), color=theme["primary"], align=PP_ALIGN.RIGHT)


def _layout_text_accent(prs, slide_item, theme, num):
    """T1. MATN: oq fon + chapda katta rangli panel"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)

    # Chapda vertikal rangli panel
    _rect(slide, 0, 0, Inches(3.6), SLIDE_H, theme["primary"])
    _rect(slide, Inches(3.6), 0, Inches(0.12), SLIDE_H, theme["accent"])

    # Panel ichida sarlavha
    _text(slide, Inches(0.5), Inches(1.2), Inches(2.8), Inches(4.5),
          slide_item.get('title', ''), size=Pt(26), bold=True, color=WHITE)

    # Katta raqam paneli pastida
    _text(slide, Inches(0.5), Inches(6.3), Inches(2.8), Inches(0.8),
          f"{num:02d}", size=Pt(40), bold=True, color=theme["accent"])

    # O'ngda matn
    _text(slide, Inches(4.3), Inches(1.0), Inches(8.3), Inches(5.2),
          slide_item.get('content', ''), size=Pt(16), color=RGBColor(0x2A, 0x2A, 0x35),
          line_spacing=1.35)

    sub = slide_item.get('subtitle', '')
    if sub:
        _rect(slide, Inches(4.3), Inches(6.35), Inches(8.3), Pt(2), theme["accent"])
        _text(slide, Inches(4.3), Inches(6.5), Inches(8.3), Inches(0.7),
              "💡 " + sub, size=Pt(12), color=theme["primary"])


def _layout_text_dark(prs, slide_item, theme, num):
    """T2. MATN: to'q fon + markazlashgan kontent, aksent chiziqlar"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme["dark"])

    # Yuqori va pastki aksent chiziqlar
    _rect(slide, 0, 0, SLIDE_W, Pt(6), theme["accent"])
    _rect(slide, 0, SLIDE_H - Pt(6), SLIDE_W, Pt(6), theme["accent"])

    # Sarlavha
    _text(slide, Inches(1.2), Inches(0.7), Inches(10.9), Inches(1.2),
          slide_item.get('title', ''), size=Pt(32), bold=True, color=WHITE)
    _rect(slide, Inches(1.2), Inches(1.8), Inches(2.4), Pt(4), theme["accent"])

    # Matn — yarim shaffof panel ichida
    _rect(slide, Inches(1.0), Inches(2.3), Inches(11.3), Inches(4.4),
          theme["primary"], transparency=55)
    _text(slide, Inches(1.5), Inches(2.6), Inches(10.3), Inches(3.8),
          slide_item.get('content', ''), size=Pt(16),
          color=RGBColor(0xEC, 0xEE, 0xF4), line_spacing=1.35)

    sub = slide_item.get('subtitle', '')
    if sub:
        _text(slide, Inches(1.2), Inches(6.9), Inches(10.9), Inches(0.5),
              "💡 " + sub, size=Pt(12), color=theme["accent"])

    _text(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
          str(num), size=Pt(12), color=WHITE, align=PP_ALIGN.RIGHT)


def _layout_text_split(prs, slide_item, theme, num):
    """T3. MATN: yuqorida rangli sarlavha bloki + pastda 2 ustunli matn"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0xF6, 0xF8, 0xFB))

    # Yuqori rangli blok
    _rect(slide, 0, 0, SLIDE_W, Inches(2.0), theme["primary"])
    _rect(slide, 0, Inches(2.0), SLIDE_W, Pt(5), theme["accent"])
    _text(slide, Inches(0.9), Inches(0.55), Inches(11.5), Inches(1.1),
          slide_item.get('title', ''), size=Pt(30), bold=True, color=WHITE)

    # Matnni 2 ustunga bo'lish
    content = slide_item.get('content', '')
    sentences = [s.strip() for s in content.replace('!', '.').replace('?', '.').split('.') if s.strip()]
    half = (len(sentences) + 1) // 2
    col1 = '. '.join(sentences[:half]) + ('.' if sentences[:half] else '')
    col2 = '. '.join(sentences[half:]) + ('.' if sentences[half:] else '')

    _text(slide, Inches(0.9), Inches(2.5), Inches(5.7), Inches(4.2),
          col1, size=Pt(14.5), color=RGBColor(0x2A, 0x2A, 0x35), line_spacing=1.3)

    # Ustunlar orasida vertikal chiziq
    _rect(slide, Inches(6.85), Inches(2.6), Pt(2.5), Inches(3.9), theme["accent"])

    _text(slide, Inches(7.15), Inches(2.5), Inches(5.3), Inches(4.2),
          col2, size=Pt(14.5), color=RGBColor(0x2A, 0x2A, 0x35), line_spacing=1.3)

    sub = slide_item.get('subtitle', '')
    if sub:
        _rect(slide, Inches(0.9), Inches(6.8), Inches(11.5), Inches(0.55), theme["primary"], transparency=88)
        _text(slide, Inches(1.1), Inches(6.85), Inches(11.1), Inches(0.45),
              "💡 " + sub, size=Pt(11.5), color=theme["primary"])

    _text(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
          str(num), size=Pt(12), color=theme["primary"], align=PP_ALIGN.RIGHT)


def _layout_bullets_noimg(prs, slide_item, theme, num):
    """T4. NUQTALI KARTALAR — rasmsiz, rangli sarlavha bilan"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0xF5, 0xF7, 0xFA))

    # Yuqori rangli banner (rasm o'rniga)
    _rect(slide, 0, 0, SLIDE_W, Inches(1.6), theme["dark"])
    _rect(slide, 0, Inches(1.6), SLIDE_W, Pt(5), theme["accent"])
    _text(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.0),
          slide_item.get('title', ''), size=Pt(30), bold=True, color=WHITE)

    items = slide_item.get('bullet_points') or []
    if not items:
        items = [s.strip() for s in slide_item.get('content', '').split('.') if s.strip()][:6]
    items = [it.strip().lstrip('-•* ') for it in items if it.strip()][:6]

    top0 = Inches(2.0)
    card_w = Inches(6.0)
    card_h = Inches(1.55)
    gap = Inches(0.25)
    for i, item in enumerate(items):
        col = i % 2
        row = i // 2
        left = Inches(0.5) + col * (card_w + gap)
        top = top0 + row * (card_h + gap)
        card = _rect(slide, left, top, card_w, card_h, WHITE)
        card.line.color.rgb = RGBColor(0xE0, 0xE4, 0xEA)
        card.line.width = Pt(1)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.18),
                                      top + (card_h - Inches(0.55)) / 2, Inches(0.55), Inches(0.55))
        circ.fill.solid()
        circ.fill.fore_color.rgb = theme["accent"] if i % 2 == 0 else theme["primary"]
        circ.line.fill.background()
        tfc = circ.text_frame
        pc = tfc.paragraphs[0]
        pc.alignment = PP_ALIGN.CENTER
        rc = pc.add_run()
        rc.text = str(i + 1)
        _set_font(rc, size=Pt(16), bold=True, color=WHITE)
        _text(slide, left + Inches(0.95), top + Inches(0.12), card_w - Inches(1.1),
              card_h - Inches(0.24), item, size=Pt(12.5),
              color=RGBColor(0x2A, 0x2A, 0x35), anchor='ctr')

    _text(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
          str(num), size=Pt(12), color=theme["primary"], align=PP_ALIGN.RIGHT)


def _layout_quote(prs, slide_item, theme, num):
    """T5. KATTA IQTIBOS/BAYONOT: markazda katta matn, aksent qavslar"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme["light"])

    # Katta dekorativ qo'shtirnoq
    _text(slide, Inches(0.6), Inches(0.2), Inches(2), Inches(1.6),
          '"', size=Pt(120), bold=True, color=theme["accent"])

    # Sarlavha
    _text(slide, Inches(2.2), Inches(1.0), Inches(9.0), Inches(1.0),
          slide_item.get('title', ''), size=Pt(28), bold=True,
          color=theme["primary"], align=PP_ALIGN.CENTER)

    # Markaziy panel — asosiy matn
    panel = _rect(slide, Inches(1.6), Inches(2.3), Inches(10.1), Inches(3.8), WHITE)
    panel.line.color.rgb = theme["accent"]
    panel.line.width = Pt(2)
    _text(slide, Inches(2.1), Inches(2.6), Inches(9.1), Inches(3.2),
          slide_item.get('content', ''), size=Pt(16),
          color=RGBColor(0x2A, 0x2A, 0x35), align=PP_ALIGN.CENTER,
          anchor='ctr', line_spacing=1.35)

    sub = slide_item.get('subtitle', '')
    if sub:
        _text(slide, Inches(1.6), Inches(6.4), Inches(10.1), Inches(0.6),
              "— " + sub, size=Pt(13), color=theme["primary"],
              align=PP_ALIGN.CENTER)

    _text(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
          str(num), size=Pt(12), color=theme["primary"], align=PP_ALIGN.RIGHT)


def _layout_timeline(prs, slide_item, theme, num):
    """T6. BOSQICHLAR/TIMELINE: gorizontal qadamlar chizig'i"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    _rect(slide, 0, 0, SLIDE_W, Pt(6), theme["primary"])

    _text(slide, Inches(0.8), Inches(0.45), Inches(11.7), Inches(1.0),
          slide_item.get('title', ''), size=Pt(30), bold=True, color=theme["dark"])
    _rect(slide, Inches(0.8), Inches(1.4), Inches(2.2), Pt(4), theme["accent"])

    # Matnni bosqichlarga bo'lish
    content = slide_item.get('content', '')
    items = slide_item.get('bullet_points') or []
    if not items:
        items = [s.strip() for s in content.replace('!', '.').split('.') if s.strip()][:4]
    items = [it.strip().lstrip('-•* ') for it in items if it.strip()][:4]

    if items:
        n_items = len(items)
        # Gorizontal chiziq
        line_y = Inches(2.6)
        _rect(slide, Inches(1.0), line_y, Inches(11.3), Pt(3), theme["accent"])

        step_w = Inches(11.3 / n_items)
        for i, item in enumerate(items):
            cx = Inches(1.0) + step_w * i + step_w / 2
            # Doira raqam
            circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                          cx - Inches(0.35), line_y - Inches(0.32),
                                          Inches(0.7), Inches(0.7))
            circ.fill.solid()
            circ.fill.fore_color.rgb = theme["primary"] if i % 2 == 0 else theme["accent"]
            circ.line.color.rgb = WHITE
            circ.line.width = Pt(2.5)
            tfc = circ.text_frame
            pc = tfc.paragraphs[0]
            pc.alignment = PP_ALIGN.CENTER
            rc = pc.add_run()
            rc.text = str(i + 1)
            _set_font(rc, size=Pt(18), bold=True, color=WHITE)

            # Bosqich matni (pastda karta)
            card = _rect(slide, Inches(1.0) + step_w * i + Inches(0.1), Inches(3.4),
                         step_w - Inches(0.2), Inches(3.2), RGBColor(0xF6, 0xF8, 0xFB))
            card.line.color.rgb = RGBColor(0xE0, 0xE4, 0xEA)
            card.line.width = Pt(1)
            _text(slide, Inches(1.0) + step_w * i + Inches(0.25), Inches(3.6),
                  step_w - Inches(0.5), Inches(2.8),
                  item, size=Pt(12), color=RGBColor(0x2A, 0x2A, 0x35), line_spacing=1.25)

    _text(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
          str(num), size=Pt(12), color=theme["primary"], align=PP_ALIGN.RIGHT)


def _layout_big_number(prs, slide_item, theme, num):
    """T7. KATTA RAQAM/FAKT: chapda ulkan raqam paneli + o'ngda matn"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme["dark"])

    # Chap panel — katta raqam
    _rect(slide, 0, 0, Inches(4.6), SLIDE_H, theme["primary"])
    _rect(slide, Inches(4.6), 0, Inches(0.12), SLIDE_H, theme["accent"])

    # Subtitle'dan raqam ajratish (masalan "50 million foydalanuvchi")
    sub = slide_item.get('subtitle', '')
    m = re.search(r'(\d[\d\s.,%]*)', sub)
    big = m.group(1).strip()[:8] if m else f"{num:02d}"
    _text(slide, Inches(0.4), Inches(2.2), Inches(3.9), Inches(1.8),
          big, size=Pt(64), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if sub:
        _text(slide, Inches(0.4), Inches(4.2), Inches(3.9), Inches(1.6),
              sub, size=Pt(13), color=RGBColor(0xE0, 0xE6, 0xF0),
              align=PP_ALIGN.CENTER)

    # O'ng tomonda sarlavha + matn
    _text(slide, Inches(5.2), Inches(0.8), Inches(7.4), Inches(1.2),
          slide_item.get('title', ''), size=Pt(27), bold=True, color=WHITE)
    _rect(slide, Inches(5.2), Inches(1.9), Inches(1.8), Pt(4), theme["accent"])
    _text(slide, Inches(5.2), Inches(2.3), Inches(7.5), Inches(4.6),
          slide_item.get('content', ''), size=Pt(15),
          color=RGBColor(0xE2, 0xE6, 0xEE), line_spacing=1.3)

    _text(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
          str(num), size=Pt(12), color=WHITE, align=PP_ALIGN.RIGHT)


def _layout_conclusion(prs, slide_item, theme, img_raw):
    """7. XULOSA: to'liq rasm + markazda matn"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = _prepare_image(img_raw, 1280, 720, darken=0.6, blur=2.0)
    _full_bg_image(slide, bg)

    _rect(slide, Inches(5.85), Inches(1.15), Inches(1.6), Pt(6), theme["accent"])
    _text(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(1.1),
          slide_item.get('title', 'Xulosa'), size=Pt(36), bold=True,
          color=WHITE, align=PP_ALIGN.CENTER)
    _text(slide, Inches(1.8), Inches(2.9), Inches(9.7), Inches(3.6),
          slide_item.get('content', ''), size=Pt(17),
          color=RGBColor(0xEC, 0xEC, 0xF2), align=PP_ALIGN.CENTER, line_spacing=1.3)


def _layout_thanks(prs, theme, img_raw, pres_title):
    """8. RAHMAT slaydi"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = _prepare_image(img_raw, 1280, 720, darken=0.65, blur=3.0)
    _full_bg_image(slide, bg)

    _text(slide, Inches(1.5), Inches(2.6), Inches(10.3), Inches(1.4),
          "E'tiboringiz uchun rahmat!", size=Pt(44), bold=True,
          color=WHITE, align=PP_ALIGN.CENTER)
    _rect(slide, Inches(5.6), Inches(4.15), Inches(2.1), Pt(5), theme["accent"])
    _text(slide, Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.9),
          pres_title, size=Pt(18), color=RGBColor(0xD8, 0xD8, 0xE2),
          align=PP_ALIGN.CENTER)


# ==================== ASOSIY GENERATOR ====================

def generate_professional_pptx(slide_data: dict, image_prompts: list[str] = None,
                                progress_callback=None, design: str = None) -> bytes:
    """
    Professional PPTX yaratish.

    design: foydalanuvchi tanlagan fan dizayni (SUBJECT_THEMES kaliti)
            bo'sh bo'lsa — AI tanlagan ranglar ishlatiladi

    slide_data = {
      "title": "...",
      "theme_colors": {"primary": "#..", "accent": "#..", "dark": "#..", "light": "#.."},
      "slides": [ {"type": "...", "title": "...", "content": "...",
                   "subtitle": "...", "bullet_points": [], "image_prompt": "..."} ]
    }
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Foydalanuvchi dizayn tanlagan bo'lsa — u ustuvor, aks holda AI ranglari
    if design and design in SUBJECT_THEMES:
        theme = build_theme(SUBJECT_THEMES[design])
        log.info(f"Dizayn tanlandi: {design}")
    else:
        theme = build_theme(slide_data.get('theme_colors'))
    slides = slide_data.get('slides', [])
    pres_title = slide_data.get('title', 'Taqdimot')

    # ---- Rasmlarni PARALLEL yuklash (tezlik uchun) ----
    def _load_one(args):
        i, s = args
        prompt = s.get('image_prompt') or (image_prompts[i] if image_prompts and i < len(image_prompts) else pres_title)
        full_prompt = (f"{prompt}, ultra realistic professional stock photography, 4k quality, "
                       f"sharp focus, cinematic lighting, detailed, "
                       f"no text, no watermark, no cartoon, no illustration")
        raw = _fetch_ai_image(full_prompt, seed=1000 + i * 7)
        if raw is None:
            log.warning(f"Rasm (slayd {i+1}) yuklanmadi")
        else:
            log.info(f"Rasm (slayd {i+1}) tayyor")
        return i, raw

    # ===== QAYSI SLAYDLARGA RASM KERAK? =====
    # Titul (0) va Xulosa (oxirgi) — har doim rasm foni
    # O'rta slaydlarning 1/4 qismi rasmli (har 4-varaq), qolganlari matn dizaynida
    n = len(slides)
    image_indices = {0}
    if n > 2:
        image_indices.add(n - 1)
        # har 4-o'rta slaydga rasm: 2, 6, 10, 14... (reja slaydi rasmsiz qoladi)
        for i in range(1, n - 1):
            if (i - 2) % 4 == 0 and i >= 2:
                image_indices.add(i)

    images = [None] * n

    if design and design in SUBJECT_IMAGES:
        # ===== FAN DIZAYNI: preview'dagi BIR XIL stok rasmlar ishlatiladi =====
        subject_urls = SUBJECT_IMAGES[design]
        url_cache = {}

        def _load_subject(args):
            pos, i = args
            url = subject_urls[pos % len(subject_urls)]
            if url in url_cache:
                return i, url_cache[url]
            raw = _fetch_url_image(url)
            url_cache[url] = raw
            return i, raw

        tasks = [(pos, i) for pos, i in enumerate(sorted(image_indices))]
        with ThreadPoolExecutor(max_workers=4) as pool:
            for i, raw in pool.map(_load_subject, tasks):
                images[i] = raw
        log.info(f"Fan dizayni '{design}': stok rasmlar yuklandi")
    else:
        # ===== AI REJIMI: mavzuga mos AI rasmlar =====
        to_load = [(i, slides[i]) for i in sorted(image_indices)]
        with ThreadPoolExecutor(max_workers=3) as pool:
            for i, raw in pool.map(_load_one, to_load):
                images[i] = raw

    # Yuklanmagan (kerakli) rasmlarga gradient zaxira
    for i in image_indices:
        if images[i] is None:
            images[i] = _fallback_gradient(1024, 576, theme)

    real_count = sum(1 for i in image_indices if images[i] is not None)
    log.info(f"Rasmlar: {real_count}/{len(image_indices)} ta yuklandi (jami {n} slaydning 1/3 qismi)")

    # Rahmat slaydi uchun rasm — birinchi rasmni ishlatamiz
    thanks_img = images[0] if images else _fallback_gradient(1280, 720, theme)

    # ---- Layout rotatsiyasi: hech qaysi ketma-ket slayd bir xil emas ----
    content_layouts = [
        _layout_image_right,
        _layout_full_image,
        _layout_image_left,
        _layout_bullets,
        _layout_split_diagonal,
    ]

    author = slide_data.get('author', '')

    img_layout_i = 0   # rasmli layoutlar hisoblagichi
    txt_layout_i = 0   # matnli layoutlar hisoblagichi

    for idx, item in enumerate(slides):
        img = images[idx]
        has_bullets = item.get('bullet_points') and len(item['bullet_points']) >= 3

        if idx == 0:
            if author:
                item = dict(item)
                item['author'] = author
            _layout_title(prs, item, theme, img, pres_title)
        elif idx == len(slides) - 1 and len(slides) > 2:
            _layout_conclusion(prs, item, theme, img)
        elif img is not None:
            # ===== RASMLI slayd (har 3-slayd) =====
            if has_bullets:
                item = dict(item)
                item['_img2'] = None
                _layout_bullets(prs, item, theme, img, idx + 1)
            else:
                img_layouts = [_layout_image_right, _layout_full_image,
                               _layout_image_left, _layout_split_diagonal]
                layout = img_layouts[img_layout_i % len(img_layouts)]
                img_layout_i += 1
                if layout is _layout_full_image:
                    layout(prs, item, theme, img, idx + 1)
                else:
                    layout(prs, item, theme, img, idx + 1)
        else:
            # ===== MATNLI slayd (rasmsiz, chiroyli dizayn) =====
            if has_bullets:
                _layout_bullets_noimg(prs, item, theme, idx + 1)
            else:
                # 6 xil matn layouti rotatsiyada (jami 15 xil varaq dizayni)
                txt_layouts = [_layout_text_accent, _layout_text_dark,
                               _layout_text_split, _layout_quote,
                               _layout_big_number, _layout_timeline]
                layout = txt_layouts[txt_layout_i % len(txt_layouts)]
                txt_layout_i += 1
                layout(prs, item, theme, idx + 1)

    _layout_thanks(prs, theme, thanks_img, pres_title)

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.getvalue()
